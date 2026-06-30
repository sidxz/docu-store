import contextlib
import io
import os
import tempfile
from pathlib import Path

import pytest

from application.ports.blob_store import StoredBlob
from infrastructure.file_services.libreoffice_converter import (
    LibreOfficeConverter,
    find_soffice,
)


class FakeBlob:
    """Blob store with on-disk source files (so get_file yields a real path)."""

    def __init__(self):
        self.sources: dict[str, bytes] = {}
        self.puts: dict[str, bytes] = {}

    def put_source(self, key, data):
        self.sources[key] = data

    @contextlib.contextmanager
    def get_file(self, key):
        d = tempfile.mkdtemp()
        p = os.path.join(d, os.path.basename(key))
        Path(p).write_bytes(self.sources[key])
        yield Path(p)

    def put_stream(self, key, stream, *, mime_type=None):
        self.puts[key] = stream.read()
        return StoredBlob(key=key, size_bytes=len(self.puts[key]), sha256="x", mime_type=mime_type)


class _FakeSoffice(LibreOfficeConverter):
    """Overrides the one external boundary: simulate LibreOffice writing <stem>.pdf."""

    def _run(self, src_path, outdir):
        (Path(outdir) / f"{Path(src_path).stem}.pdf").write_bytes(b"%PDF-1.4 fake")


def test_convert_to_pdf_stores_pdf_at_dest_key():
    blob = FakeBlob()
    blob.put_source("artifacts/x/source.pptx", b"pptx-bytes")
    conv = _FakeSoffice(blob_store=blob, soffice_bin="/bin/true")

    conv.convert_to_pdf("artifacts/x/source.pptx", "artifacts/x/derived/render.pdf")

    assert blob.puts["artifacts/x/derived/render.pdf"].startswith(b"%PDF")


def test_convert_raises_when_no_pdf_produced():
    class _NoOutput(LibreOfficeConverter):
        def _run(self, src_path, outdir):
            pass  # produce nothing

    blob = FakeBlob()
    blob.put_source("artifacts/x/source.pptx", b"pptx-bytes")
    conv = _NoOutput(blob_store=blob, soffice_bin="/bin/true")

    with pytest.raises(RuntimeError):
        conv.convert_to_pdf("artifacts/x/source.pptx", "artifacts/x/derived/render.pdf")


@pytest.mark.skipif(find_soffice() is None, reason="LibreOffice not installed")
def test_real_soffice_converts_pptx():
    from pptx import Presentation

    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[1])
    s.shapes.title.text = "Hello"
    buf = io.BytesIO()
    prs.save(buf)

    blob = FakeBlob()
    blob.put_source("artifacts/x/source.pptx", buf.getvalue())
    conv = LibreOfficeConverter(blob_store=blob)

    conv.convert_to_pdf("artifacts/x/source.pptx", "artifacts/x/derived/render.pdf")

    assert blob.puts["artifacts/x/derived/render.pdf"].startswith(b"%PDF")
